#!/usr/bin/env python3
"""
Generate a comprehensive PowerPoint presentation for the Bachelor Thesis:
"Smart Task & Quality Management System (TaskFlow)" by Ahmed Hassan Abdelbarr
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
import os

# ── Paths ──────────────────────────────────────────────────────────────────
BASE = "/home/ahmed-abdelbarr/Github/Task_Flow"
FIGS = os.path.join(BASE, "Bachelor_Thesis", "Figures")
OUTPUT = os.path.join(BASE, "TaskFlow_Thesis_Presentation.pptx")

# ── Color Palette ──────────────────────────────────────────────────────────
DARK_NAVY   = RGBColor(0x0F, 0x17, 0x2A)   # Background / footer bars
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)   # Headings, accent lines
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF1, 0xF5, 0xF9)   # Slide background (near-white)
MED_GRAY    = RGBColor(0x94, 0xA3, 0xB8)   # Subtle text
DARK_TEXT    = RGBColor(0x1E, 0x29, 0x3B)   # Body text
GREEN_OK    = RGBColor(0x10, 0xB9, 0x81)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
RED_ACCENT  = RGBColor(0xEF, 0x44, 0x44)
TEAL        = RGBColor(0x14, 0xB8, 0xA6)
PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)

# ── Helpers ────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLD_W = prs.slide_width
SLD_H = prs.slide_height

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text="", font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(4)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return tb

def add_para(text_frame, text, font_size=16, bold=False, color=DARK_TEXT,
             alignment=PP_ALIGN.LEFT, font_name="Calibri", space_before=0,
             space_after=6, bullet=False, level=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.level = level
    if bullet:
        p.level = level
    return p

def add_bullet_text(slide, left, top, width, height, items, font_size=15,
                    color=DARK_TEXT, font_name="Calibri", bold_first=False,
                    line_spacing=1.3):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.line_spacing = Pt(font_size * line_spacing)
        if bold_first and i == 0:
            p.font.bold = True
    return tb

def add_image_safe(slide, filename, left, top, width=None, height=None):
    """Add an image if file exists, returns shape or None."""
    path = os.path.join(FIGS, filename)
    if not os.path.exists(path):
        # Try alternate extensions
        for ext in ['.jpg', '.jpeg', '.png', '.gif']:
            base = os.path.splitext(filename)[0]
            alt = os.path.join(FIGS, base + ext)
            if os.path.exists(alt):
                path = alt
                break
        else:
            # Placeholder rectangle
            shape = add_rect(slide, left, top, width or Inches(3), height or Inches(2),
                             MED_GRAY)
            tb = add_text_box(slide, left, top, width or Inches(3), height or Inches(2),
                              f"[Image: {filename}]", font_size=12, color=WHITE,
                              alignment=PP_ALIGN.CENTER)
            return shape
    if width and height:
        return slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        return slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        return slide.shapes.add_picture(path, left, top, height=height)
    else:
        return slide.shapes.add_picture(path, left, top)

def add_section_header(slide, number, title, subtitle=None):
    """Dark sidebar with section number and title."""
    # Full-width dark bar at top
    add_rect(slide, 0, 0, SLD_W, Inches(1.6), DARK_NAVY)
    # Accent line
    add_rect(slide, Inches(0.6), Inches(1.5), Inches(1.2), Pt(4), ACCENT_BLUE)

    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(2), Inches(0.5),
                 f"0{number}" if number < 10 else str(number),
                 font_size=36, bold=True, color=ACCENT_BLUE)
    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(11), Inches(0.8),
                 title, font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(1.1), Inches(11), Inches(0.5),
                     subtitle, font_size=16, color=MED_GRAY)

def add_footer(slide, text="Smart Task & Quality Management System (TaskFlow)  |  Ahmed Hassan Abdelbarr"):
    add_rect(slide, 0, Inches(7.1), SLD_W, Inches(0.4), DARK_NAVY)
    add_text_box(slide, Inches(0.5), Inches(7.12), Inches(12), Inches(0.35),
                 text, font_size=9, color=MED_GRAY, alignment=PP_ALIGN.LEFT)


def make_table(slide, rows, cols, data, left, top, width, height,
               col_widths=None, header_color=DARK_NAVY, header_font_color=WHITE):
    """Create a styled table. data is list of lists (first row = header)."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_font_color
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_rect(sl, 0, 0, SLD_W, SLD_H, DARK_NAVY)
# Decorative accent shapes
add_rect(sl, 0, 0, Inches(0.15), SLD_H, ACCENT_BLUE)
add_rect(sl, Inches(0.15), Inches(3.0), Inches(0.08), Inches(1.8), TEAL)

# GUC Logo
add_image_safe(sl, "GUC_logo.png", Inches(0.8), Inches(0.4), height=Inches(0.9))

add_text_box(sl, Inches(0.8), Inches(1.6), Inches(11), Inches(0.4),
             "Bachelor of Science in Computer Science and Engineering",
             font_size=14, color=MED_GRAY)

add_text_box(sl, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.4),
             "Smart Task & Quality Management System",
             font_size=44, bold=True, color=WHITE)

add_text_box(sl, Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.6),
             "TASKFLOW",
             font_size=28, bold=True, color=ACCENT_BLUE)

add_text_box(sl, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.5),
             "An Offline-First Cross-Platform Desktop Application with AI-Powered Assistance",
             font_size=18, color=LIGHT_GRAY)

add_text_box(sl, Inches(0.8), Inches(5.4), Inches(5), Inches(0.4),
             "Ahmed Hassan Abdelbarr", font_size=20, bold=True, color=WHITE)
add_text_box(sl, Inches(0.8), Inches(5.85), Inches(5), Inches(0.4),
             "Supervised by: Assoc. Prof. Wael Zakaria Abdallah",
             font_size=14, color=MED_GRAY)
add_text_box(sl, Inches(0.8), Inches(6.3), Inches(5), Inches(0.4),
             "German University in Cairo  •  Faculty of Engineering and Materials Science  •  2025",
             font_size=12, color=MED_GRAY)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_footer(sl)

# Left sidebar
add_rect(sl, 0, 0, Inches(3.8), SLD_H, DARK_NAVY)
add_rect(sl, 0, 0, Inches(0.12), SLD_H, ACCENT_BLUE)
add_text_box(sl, Inches(0.5), Inches(0.5), Inches(3), Inches(0.5),
             "Agenda", font_size=32, bold=True, color=WHITE)

agenda_items = [
    ("01", "Problem Statement"),
    ("02", "Objectives"),
    ("03", "Background & Related Work"),
    ("04", "System Architecture & Design"),
    ("05", "Key Innovations"),
    ("06", "Results & Evaluation"),
    ("07", "Contributions"),
    ("08", "Future Work"),
    ("09", "Live Demo"),
    ("10", "Q&A"),
]
for i, (num, title) in enumerate(agenda_items):
    y = Inches(1.3) + Inches(i * 0.58)
    add_text_box(sl, Inches(0.6), y, Inches(0.6), Inches(0.35),
                 num, font_size=14, bold=True, color=ACCENT_BLUE)
    add_text_box(sl, Inches(1.3), y, Inches(2.5), Inches(0.35),
                 title, font_size=15, color=WHITE)

# Right side — preview text
add_text_box(sl, Inches(4.5), Inches(1.0), Inches(8), Inches(0.5),
             "Presentation Overview", font_size=26, bold=True, color=DARK_NAVY)
add_text_box(sl, Inches(4.5), Inches(1.7), Inches(8), Inches(4.5),
             "This presentation covers the design, implementation, and evaluation of TaskFlow — "
             "a cross-platform desktop task management system built on an offline-first architecture "
             "with dual-database synchronisation, multi-modal AI integration, and a rich collaborative feature set.\n\n"
             "Duration: ~15 minutes  |  15 slides  |  Live demo included",
             font_size=16, color=DARK_TEXT, line_spacing=1.5)
# Decorative line
add_rect(sl, Inches(4.5), Inches(6.5), Inches(3), Pt(3), ACCENT_BLUE)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_header(sl, 1, "Problem Statement",
                   "Why existing task management tools fall short")
add_footer(sl)

# Three problem columns
col_w = Inches(3.6)
col_gap = Inches(0.4)
col_start = Inches(0.6)
col_y = Inches(2.0)
col_h = Inches(4.8)

problems = [
    ("Cloud-Centric Fragility",
     "cloud_fragile",
     ACCENT_BLUE,
     [
         "Tools like Trello, Asana, Jira, and Notion assume always-on connectivity",
         "When offline, operations are silently lost or error out",
         "Unusable on trains, planes, rural areas, or during network outages",
         "No graceful degradation — full functionality requires internet",
     ]),
    ("Fragmented Toolchains",
     "fragmented",
     RED_ACCENT,
     [
         "Users juggle separate apps for tasks, code, calendar, chat, and notifications",
         "No single cohesive desktop application integrates all of these",
         "Context-switching reduces productivity and increases cognitive load",
         "Task assignments buried in social media / messaging apps with no structure",
     ]),
    ("Shallow AI Integration",
     "shallow_ai",
     TEAL,
     [
         "AI features in existing tools are limited to basic text generation",
         "No conversational assistant aware of the user's own project data",
         "No support for code generation, document OCR, or multimodal input",
         "Missed opportunity for AI to truly augment project management",
     ]),
]

for i, (title, key, color, bullets) in enumerate(problems):
    x = col_start + i * (col_w + col_gap)
    # Card background
    card = add_rect(sl, x, col_y, col_w, col_h, WHITE)
    card.shadow.inherit = False
    # Top accent bar
    add_rect(sl, x, col_y, col_w, Pt(5), color)
    # Icon circle
    circle = add_rect(sl, x + Inches(1.35), col_y + Inches(0.25), Inches(0.9), Inches(0.9), color)
    # Title
    add_text_box(sl, x + Inches(0.2), col_y + Inches(1.3), col_w - Inches(0.4), Inches(0.5),
                 title, font_size=17, bold=True, color=DARK_NAVY, alignment=PP_ALIGN.CENTER)
    # Bullets
    bullet_y = col_y + Inches(2.0)
    for j, b in enumerate(bullets):
        tb = add_text_box(sl, x + Inches(0.25), bullet_y + Inches(j * 0.65),
                          col_w - Inches(0.5), Inches(0.6),
                          f"  {b}", font_size=12, color=DARK_TEXT)
        # Bullet dot
        add_rect(sl, x + Inches(0.25), bullet_y + Inches(j * 0.65) + Pt(5),
                 Pt(5), Pt(5), color)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — OBJECTIVES
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_header(sl, 2, "Objectives",
                   "Six design goals that guided the development of TaskFlow")
add_footer(sl)

objectives = [
    ("Offline-First\nOperation", "Every core feature works fully\nwithout internet connectivity.\nSQLite for local persistence.", ACCENT_BLUE),
    ("Reliable\nSynchronisation", "Transactional outbox pattern\nfor at-least-once delivery.\nUser-visible sync progress.", GREEN_OK),
    ("Multi-Mode\nAI Assistant", "Mistral API: chat, code gen\n(Codestral), and document OCR.\nSSE streaming for responsiveness.", PURPLE),
    ("Five Task\nVisualisations", "Default, Kanban, Table, Gantt,\nand Calendar views — each\noptimised for different workflows.", AMBER),
    ("Cross-Platform\nDesktop App", "Self-contained ~35MB installer\nvia Tauri v2. Embedded .NET\nbackend + React frontend.", TEAL),
    ("Full\nCollaboration", "Teams, DMs, group chats, file\nattachments, real-time SignalR\nnotifications, calendar events.", RED_ACCENT),
]

box_w = Inches(3.8)
box_h = Inches(3.0)
gap = Inches(0.3)
start_x = Inches(0.6)
start_y = Inches(2.0)

for i, (title, desc, color) in enumerate(objectives):
    col = i % 3
    row = i // 3
    x = start_x + col * (box_w + gap)
    y = start_y + row * (box_h + gap)

    card = add_rect(sl, x, y, box_w, box_h, WHITE)
    add_rect(sl, x, y, box_w, Pt(5), color)
    # Number circle
    num_shape = add_rect(sl, x + Inches(0.2), y + Inches(0.3), Inches(0.6), Inches(0.6), color)
    add_text_box(sl, x + Inches(0.2), y + Inches(0.3), Inches(0.6), Inches(0.6),
                 str(i + 1), font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(1.0), y + Inches(0.25), Inches(2.6), Inches(0.7),
                 title, font_size=16, bold=True, color=DARK_NAVY)
    add_text_box(sl, x + Inches(0.3), y + Inches(1.1), Inches(3.3), Inches(1.7),
                 desc, font_size=13, color=DARK_TEXT, line_spacing=1.4)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — BACKGROUND & RELATED WORK
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_header(sl, 3, "Background & Related Work",
                   "Survey of existing tools and technologies")
add_footer(sl)

# Comparison table
table_data = [
    ["Feature", "Trello", "Asana", "Jira", "Notion", "TaskFlow"],
    ["Offline Operation",     "✗", "✗", "✗", "Partial", "✓ Full"],
    ["Local Persistence",     "✗", "✗", "✗", "✗", "✓ SQLite"],
    ["Transparent Sync",      "✗", "✗", "✗", "✗", "✓ Outbox Pattern"],
    ["AI Assistant",          "✗", "✗", "✗", "Basic", "✓ Multi-Mode (Mistral)"],
    ["Desktop Installer",     "✗", "✗", "✗", "✓", "✓ Tauri (~35MB)"],
    ["Kanban View",           "✓", "✓", "✓", "✓", "✓"],
    ["Gantt View",            "✗", "✓", "✓", "✗", "✓"],
    ["Calendar View",         "✗", "✓", "✗", "✓", "✓"],
    ["Team Messaging",        "✗", "✓", "✗", "✓", "✓ SignalR"],
]

tbl = make_table(sl, len(table_data), len(table_data[0]), table_data,
                 Inches(0.6), Inches(2.0), Inches(12.1), Inches(3.5),
                 col_widths=[Inches(2.0), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(4.1)])

# Bottom key insight
add_rect(sl, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.0), RGBColor(0xEE, 0xF2, 0xFF))
add_text_box(sl, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.8),
             " Key Insight: TaskFlow is the only tool combining offline-first architecture, local SQLite persistence, "
             "transparent synchronisation, multi-modal AI, and a self-contained desktop installer — "
             "all in a single cohesive application.",
             font_size=14, bold=True, color=DARK_NAVY)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_header(sl, 4, "System Architecture",
                   "Three-tier desktop application with dual-database strategy")
add_footer(sl)

# Architecture diagram as shapes
# Tier 1 - Presentation
y_t1 = Inches(2.1)
h_tier = Inches(1.3)
gap_tier = Inches(0.15)

# Presentation Layer
box = add_rect(sl, Inches(0.6), y_t1, Inches(12.1), h_tier, RGBColor(0xDB, 0xE8, 0xFE))
add_rect(sl, Inches(0.6), y_t1, Inches(0.12), h_tier, ACCENT_BLUE)
add_text_box(sl, Inches(0.9), y_t1 + Inches(0.1), Inches(3), Inches(0.3),
             "PRESENTATION LAYER", font_size=12, bold=True, color=ACCENT_BLUE)
add_text_box(sl, Inches(0.9), y_t1 + Inches(0.4), Inches(5), Inches(0.7),
             "React 19 + TypeScript SPA\nREST HTTP + SignalR WebSocket", font_size=14, color=DARK_TEXT)

# Down arrow
add_text_box(sl, Inches(6.3), y_t1 + h_tier - Inches(0.1), Inches(0.6), Inches(0.4),
             "▼", font_size=18, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Tier 2 - Application Logic
y_t2 = y_t1 + h_tier + gap_tier + Inches(0.15)
box2 = add_rect(sl, Inches(0.6), y_t2, Inches(12.1), h_tier, RGBColor(0xD1, 0xFA, 0xE5))
add_rect(sl, Inches(0.6), y_t2, Inches(0.12), h_tier, GREEN_OK)
add_text_box(sl, Inches(0.9), y_t2 + Inches(0.1), Inches(3), Inches(0.3),
             "APPLICATION LOGIC LAYER", font_size=12, bold=True, color=GREEN_OK)
add_text_box(sl, Inches(0.9), y_t2 + Inches(0.4), Inches(5), Inches(0.7),
             "ASP.NET Core (.NET 10)  •  C#\nControllers → Services → Repositories", font_size=14, color=DARK_TEXT)
# Key services badges
svcs = ["Offline Sync", "SignalR Hub", "AI Gateway", "Auth", "Notification"]
for j, s in enumerate(svcs):
    bx = add_rect(sl, Inches(6.5) + Inches(j * 1.15), y_t2 + Inches(0.4),
                  Inches(1.05), Inches(0.5), GREEN_OK)
    add_text_box(sl, Inches(6.5) + Inches(j * 1.15), y_t2 + Inches(0.4),
                 Inches(1.05), Inches(0.5), s, font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER)

# Down arrow
add_text_box(sl, Inches(6.3), y_t2 + h_tier - Inches(0.1), Inches(0.6), Inches(0.4),
             "▼", font_size=18, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Tier 3 - Data Layer
y_t3 = y_t2 + h_tier + gap_tier + Inches(0.15)
box3 = add_rect(sl, Inches(0.6), y_t3, Inches(12.1), h_tier, RGBColor(0xFE, 0xF3, 0xC7))
add_rect(sl, Inches(0.6), y_t3, Inches(0.12), h_tier, AMBER)
add_text_box(sl, Inches(0.9), y_t3 + Inches(0.1), Inches(3), Inches(0.3),
             "DATA LAYER", font_size=12, bold=True, color=AMBER)

# SQLite block
sql_box = add_rect(sl, Inches(0.9), y_t3 + Inches(0.4), Inches(3.5), Inches(0.7),
                   RGBColor(0xFF, 0xF3, 0xCD))
add_text_box(sl, Inches(1.0), y_t3 + Inches(0.42), Inches(3.3), Inches(0.65),
             "SQLite  •  18 Entities  •  ACID Compliant\nPrimary local store — zero network dependency",
             font_size=12, color=DARK_TEXT)

# Sync arrow
add_text_box(sl, Inches(4.7), y_t3 + Inches(0.45), Inches(1.5), Inches(0.6),
             "◄── Sync ──►", font_size=12, bold=True, color=DARK_NAVY, alignment=PP_ALIGN.CENTER)

# MongoDB block
mongo_box = add_rect(sl, Inches(6.2), y_t3 + Inches(0.4), Inches(3.5), Inches(0.7),
                     RGBColor(0xFF, 0xF3, 0xCD))
add_text_box(sl, Inches(6.3), y_t3 + Inches(0.42), Inches(3.3), Inches(0.65),
             "MongoDB  •  15 Collections  •  Document Store\nCloud relay — cross-device synchronisation",
             font_size=12, color=DARK_TEXT)

# Tauri shell — side label
add_rect(sl, Inches(10.0), Inches(2.1), Inches(2.7), Inches(4.7), DARK_NAVY)
add_text_box(sl, Inches(10.1), Inches(2.2), Inches(2.5), Inches(0.3),
             "DESKTOP SHELL", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(sl, Inches(10.1), Inches(2.6), Inches(2.5), Inches(3.5),
             "Tauri v2 (Rust)\n\n• OS-native WebView2\n• Sidecar .NET backend\n• Stdout port negotiation\n"
             "• Splashscreen → Main\n  window lifecycle\n• 9 IPC commands\n• CSP security policy\n"
             "• NSIS / MSI installer\n  (~35 MB self-contained)",
             font_size=12, color=WHITE, line_spacing=1.3)


# ── Include ER diagram images ─────────────────────────────────────────
add_text_box(sl, Inches(0.8), Inches(5.5), Inches(4), Inches(0.3),
             "Data Model — Entity Relationship Diagrams:",
             font_size=13, bold=True, color=DARK_TEXT)
add_image_safe(sl, "erd_sqlite.png", Inches(0.8), Inches(5.9), height=Inches(1.0))
add_image_safe(sl, "erd_mongo.png", Inches(4.2), Inches(5.9), height=Inches(1.0))


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — OFFLINE-FIRST: TRANSACTIONAL OUTBOX PATTERN
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_header(sl, 5, "Key Innovation: Offline-First Synchronisation",
                   "Transactional Outbox Pattern adapted for desktop .NET")
add_footer(sl)

# Flow diagram
# Step 1: Local Write
y_flow = Inches(2.1)
step_h = Inches(2.8)
step_w = Inches(3.6)

# Local SQLite box
box_l = add_rect(sl, Inches(0.5), y_flow, step_w, step_h, WHITE)
add_rect(sl, Inches(0.5), y_flow, step_w, Pt(5), ACCENT_BLUE)
add_text_box(sl, Inches(0.7), y_flow + Inches(0.2), Inches(3.2), Inches(0.3),
             "1. LOCAL SQLITE TRANSACTION", font_size=13, bold=True, color=ACCENT_BLUE)
add_text_box(sl, Inches(0.7), y_flow + Inches(0.65), Inches(3.2), Inches(1.8),
             "Within a single atomic transaction:\n\n"
             "  ├── Write entity changes\n"
             "  │   (TaskItem, Message, etc.)\n"
             "  │\n"
             "  └── Write SyncOutboxEntry\n"
             "      { EntityType, SyncId,\n"
             "        Operation, Payload,\n"
             "        Status: Pending }\n\n"
             "  ✓ ACID guarantee — both or neither",
             font_size=12, color=DARK_TEXT, line_spacing=1.3)

# Arrow
add_text_box(sl, Inches(4.2), y_flow + Inches(1.0), Inches(0.6), Inches(0.4),
             "▶", font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Background Worker box
box_wk = add_rect(sl, Inches(4.8), y_flow, step_w, step_h, WHITE)
add_rect(sl, Inches(4.8), y_flow, step_w, Pt(5), GREEN_OK)
add_text_box(sl, Inches(5.0), y_flow + Inches(0.2), Inches(3.2), Inches(0.3),
             "2. OFFLINESYNCSERVICE", font_size=13, bold=True, color=GREEN_OK)
add_text_box(sl, Inches(5.0), y_flow + Inches(0.65), Inches(3.2), Inches(1.8),
             "Background worker polls every\nfew seconds:\n\n"
             "  ├── Reads Pending outbox\n"
             "  │   entries (FIFO order)\n"
             "  │\n"
             "  ├── Marks → Processing\n"
             "  │\n"
             "  └── Uses ConnectivityService\n"
             "      as circuit breaker\n\n"
             "  ↻ Retries Failed entries",
             font_size=12, color=DARK_TEXT, line_spacing=1.3)

# Arrow
add_text_box(sl, Inches(8.5), y_flow + Inches(1.0), Inches(0.6), Inches(0.4),
             "▶", font_size=24, color=GREEN_OK, alignment=PP_ALIGN.CENTER)

# MongoDB box
box_r = add_rect(sl, Inches(9.1), y_flow, step_w + Inches(0.3), step_h, WHITE)
add_rect(sl, Inches(9.1), y_flow, step_w + Inches(0.3), Pt(5), AMBER)
add_text_box(sl, Inches(9.3), y_flow + Inches(0.2), Inches(3.4), Inches(0.3),
             "3. MONGODB RELAY", font_size=13, bold=True, color=AMBER)
add_text_box(sl, Inches(9.3), y_flow + Inches(0.65), Inches(3.4), Inches(1.8),
             "HTTP PUT to MongoDB:\n\n"
             "  ├── Idempotent upsert\n"
             "  │   (SyncId as _id)\n"
             "  │\n"
             "  ├── On success → Synced\n"
             "  │\n"
             "  └── On failure → Failed\n"
             "      (retry with backoff)\n\n"
             "  State machine:\n"
             "  Pending → Processing → Synced",
             font_size=12, color=DARK_TEXT, line_spacing=1.3)

# Bottom: connectivity states + UI
add_rect(sl, Inches(0.5), Inches(5.2), Inches(12.2), Inches(1.6), RGBColor(0xF1, 0xF5, 0xF9))

add_text_box(sl, Inches(0.7), Inches(5.3), Inches(3), Inches(0.3),
             "Connectivity States:", font_size=14, bold=True, color=DARK_NAVY)

states = [
    ("Online", GREEN_OK, "Normal operation\nOutbox replayed immediately"),
    ("Manually Offline", AMBER, "User-toggled\nOperations queued locally"),
    ("Disconnected", RED_ACCENT, "No connection\nQueue grows; auto-replay on reconnect"),
]
for i, (state, color, desc) in enumerate(states):
    x = Inches(0.7) + Inches(i * 4.1)
    add_rect(sl, x, Inches(5.7), Inches(0.7), Inches(0.3), color)
    add_text_box(sl, x, Inches(5.7), Inches(0.7), Inches(0.3),
                 state, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.8), Inches(5.7), Inches(3.0), Inches(0.8),
                 desc, font_size=11, color=DARK_TEXT, line_spacing=1.2)

add_text_box(sl, Inches(0.7), Inches(6.5), Inches(11), Inches(0.3),
             "UI: ConnectivityBar shows real-time status • Pending operation count • Colour-coded indicators • Animated sync progress",
             font_size=11, color=MED_GRAY)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — AI INTEGRATION
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_header(sl, 5, "Key Innovation: Multi-Modal AI Assistant",
                   "Mistral API integration with three specialised modes")
add_footer(sl)

# Three mode cards
modes = [
    ("💬  General Chat", "mistral-small-latest", ACCENT_BLUE,
     "Conversational Q&A\nBrainstorming ideas\nHelp & guidance\nProject advice",
     "res_ai_chatbot_empty.jpg"),
    ("💻  Code Generation", "codestral-latest", PURPLE,
     "Generate code snippets\nExplain code logic\nRefactor existing code\nDebug assistance",
     "res_ai_coder.jpg"),
    ("📄  Document OCR", "mistral-ocr-latest", TEAL,
     "Extract text from images\nAnalyse documents\nProcess uploaded files\nMultimodal understanding",
     "res_ai_smart_autofill.png.png"),
]

for i, (title, model, color, features, img) in enumerate(modes):
    x = Inches(0.5) + Inches(i * 4.2)
    w = Inches(3.9)

    card = add_rect(sl, x, Inches(2.1), w, Inches(4.2), WHITE)
    add_rect(sl, x, Inches(2.1), w, Pt(5), color)

    add_text_box(sl, x + Inches(0.2), Inches(2.3), w - Inches(0.4), Inches(0.4),
                 title, font_size=17, bold=True, color=DARK_NAVY)
    add_text_box(sl, x + Inches(0.2), Inches(2.8), w - Inches(0.4), Inches(0.3),
                 f"Model: {model}", font_size=11, color=MED_GRAY)
    add_text_box(sl, x + Inches(0.2), Inches(3.3), w - Inches(0.4), Inches(1.6),
                 features, font_size=13, color=DARK_TEXT, line_spacing=1.5)

    # Screenshot
    add_image_safe(sl, img, x + Inches(0.2), Inches(4.9), width=w - Inches(0.4))

# Bottom technical details
add_rect(sl, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.55), RGBColor(0xEE, 0xF2, 0xFF))
add_text_box(sl, Inches(0.7), Inches(6.52), Inches(11.9), Inches(0.5),
             "Technical: SSE (Server-Sent Events) streaming from React → ASP.NET Core → Mistral API • Sliding window context (last 20 messages) • "
             "AI-generated conversation titles • Persisted conversation history in SQLite",
             font_size=11, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — FIVE TASK VISUALISATIONS
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_header(sl, 5, "Key Innovation: Five Task Visualisation Modes",
                   "Multiple perspectives on project data, each optimised for different workflows")
add_footer(sl)

views = [
    ("Default", "Smart temporal grouping\n(Overdue, Today, This Week,\nNext Week, Future, No Date)", "res_default_view.jpg"),
    ("Kanban", "Drag-and-drop columns\n(To Do, In Progress, Review,\nCompleted, Overdue)", "res_kanban_view.jpg"),
    ("Table", "Sortable, filterable grid\nwith inline editing\nand column customisation", "res_table_view.jpg"),
    ("Gantt", "Timeline-based view\nfor project scheduling\nand dependency tracking", "res_gantt_view.jpg"),
    ("Calendar", "Weekly time-grid\nwith drag-to-schedule\nand event overlays", "res_calendar_view.jpg"),
]

for i, (name, desc, img) in enumerate(views):
    x = Inches(0.3) + Inches(i * 2.6)
    w = Inches(2.4)

    add_text_box(sl, x, Inches(2.0), w, Inches(0.35),
                 f"  {i+1}. {name}", font_size=14, bold=True, color=DARK_NAVY)
    add_image_safe(sl, img, x, Inches(2.4), width=w, height=Inches(1.6))
    add_text_box(sl, x, Inches(4.1), w, Inches(1.2),
                 desc, font_size=11, color=DARK_TEXT, line_spacing=1.3)

# Bottom: supplementary screenshots
add_text_box(sl, Inches(0.5), Inches(5.5), Inches(6), Inches(0.3),
             "Additional Interface Screenshots:", font_size=14, bold=True, color=DARK_NAVY)

extra_imgs = [
    ("Dashboard", "res_dashboard_populated.jpg"),
    ("Task Editor", "res_task_editor_empty.png.png"),
]
for i, (label, img) in enumerate(extra_imgs):
    x = Inches(0.5) + Inches(i * 4.0)
    add_text_box(sl, x, Inches(5.85), Inches(2), Inches(0.25), label, font_size=11, bold=True, color=DARK_TEXT)
    add_image_safe(sl, img, x, Inches(6.1), width=Inches(3.5), height=Inches(1.0))


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — RESULTS: REQUIREMENTS MATRIX
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_header(sl, 6, "Results: Requirements Fulfilment",
                   "All 11 functional and 5 non-functional requirements fully met")
add_footer(sl)

# FR Table
fr_data = [
    ["#", "Functional Requirement", "Status"],
    ["FR-01", "User Authentication (register, login, password reset)", "✓ Fully Met"],
    ["FR-02", "Task Management (CRUD, priorities, statuses, starring, assignment)", "✓ Fully Met"],
    ["FR-03", "Multiple Task Views (Default, Kanban, Table, Gantt, Calendar)", "✓ Fully Met"],
    ["FR-04", "Projects & Teams (membership, MongoDB-relayed invitations)", "✓ Fully Met"],
    ["FR-05", "Messaging (DMs, group chats, file attachments up to 20 MB)", "✓ Fully Met"],
    ["FR-06", "Notifications (SignalR real-time, email reminders, 17 types)", "✓ Fully Met"],
    ["FR-07", "Calendar Events (CRUD, colour, meeting links, weekly grid)", "✓ Fully Met"],
    ["FR-08", "AI Chatbot (3 modes, SSE streaming, persisted history)", "✓ Fully Met"],
    ["FR-09", "Offline Operation (outbox queue, auto-replay, BulkSync)", "✓ Fully Met"],
    ["FR-10", "Dashboard (stats, task trend, due this week, recent activity)", "✓ Fully Met"],
    ["FR-11", "Settings (profile, password, notification preferences)", "✓ Fully Met"],
]

tbl_fr = make_table(sl, len(fr_data), 3, fr_data,
                    Inches(0.5), Inches(1.9), Inches(7.5), Inches(4.9),
                    col_widths=[Inches(0.7), Inches(4.8), Inches(2.0)])

# NFR Summary on the right
add_rect(sl, Inches(8.3), Inches(1.9), Inches(4.5), Inches(4.9), WHITE)
add_rect(sl, Inches(8.3), Inches(1.9), Inches(4.5), Pt(5), TEAL)
add_text_box(sl, Inches(8.5), Inches(2.1), Inches(4.1), Inches(0.3),
             "NON-FUNCTIONAL REQ.", font_size=14, bold=True, color=TEAL)
nfr_items = [
    "NFR-01: Offline-first architecture  ✓",
    "NFR-02: Cross-platform (Windows)    ✓",
    "NFR-03: Real-time collaboration     ✓",
    "NFR-04: Performance (<200ms ops)    ✓",
    "NFR-05: Security (bcrypt, CSP)      ✓",
]
for i, item in enumerate(nfr_items):
    add_text_box(sl, Inches(8.5), Inches(2.6) + Inches(i * 0.7), Inches(4.1), Inches(0.5),
                 item, font_size=13, color=DARK_TEXT, line_spacing=1.3)

# Summary banner
add_rect(sl, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5), GREEN_OK)
add_text_box(sl, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
             "✓  ALL 11 FUNCTIONAL AND 5 NON-FUNCTIONAL REQUIREMENTS SUCCESSFULLY IMPLEMENTED  ✓",
             font_size=15, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — PERFORMANCE BENCHMARKS
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_header(sl, 6, "Results: Performance Benchmarks",
                   "Measured on Windows 11, Intel Core i7, 16GB RAM, NVMe SSD")
add_footer(sl)

perf_data = [
    ["Operation", "Measured Time", "Threshold", "Result"],
    ["Application cold start (to UI ready)", "~1,500 ms", "—", "✓ Acceptable"],
    ["Login (local identity + MongoDB pull)", "<500 ms", "2,000 ms", "✓ Pass"],
    ["Task list load (100 tasks, SQLite)", "<100 ms", "200 ms", "✓ Pass"],
    ["Task create (SQLite write + SignalR)", "<80 ms", "200 ms", "✓ Pass"],
    ["First AI token rendered (General)", "~400 ms", "2,000 ms", "✓ Pass"],
    ["Offline-to-online sync replay (20 entries)", "~1,800 ms", "—", "✓ Acceptable"],
    ["Direct message delivery (sender → recipient)", "<150 ms", "200 ms", "✓ Pass"],
    ["Dashboard load", "<200 ms", "200 ms", "✓ Pass"],
    ["Kanban drag-and-drop update", "<100 ms", "200 ms", "✓ Pass"],
]

tbl_perf = make_table(sl, len(perf_data), 4, perf_data,
                      Inches(0.6), Inches(1.9), Inches(12.1), Inches(4.2),
                      col_widths=[Inches(4.5), Inches(2.5), Inches(2.5), Inches(2.6)])

# Key insight box
add_rect(sl, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.7), RGBColor(0xEE, 0xF2, 0xFF))
add_text_box(sl, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.6),
             " All interactive operations satisfy the 200ms responsiveness threshold (NFR-04). "
             "Cold start and AI first-token are dominated by process initialisation and network round-trip respectively.",
             font_size=13, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — OFFLINE SYNC VERIFICATION
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_header(sl, 6, "Results: Offline Synchronisation Verification",
                   "End-to-end test of the outbox pattern")
add_footer(sl)

# Test procedure and results side by side
# Left: Procedure
add_rect(sl, Inches(0.5), Inches(2.1), Inches(6.0), Inches(4.5), WHITE)
add_rect(sl, Inches(0.5), Inches(2.1), Inches(6.0), Pt(5), ACCENT_BLUE)
add_text_box(sl, Inches(0.7), Inches(2.3), Inches(5.5), Inches(0.3),
             "TEST PROCEDURE", font_size=16, bold=True, color=ACCENT_BLUE)

steps = [
    "1. User toggles to Manual Offline mode",
    "2. User performs 20 task operations (create, update, delete)",
    "3. All 20 SyncOutboxEntry rows created with Pending status",
    "4. No data reaches MongoDB during offline period",
    "5. User reconnects to network",
    "6. OfflineSyncService replays entries in FIFO order",
    "7. All 20 entries transition: Pending → Processing → Synced",
    "8. Tasks visible in MongoDB and on other devices via BulkSync",
]
for i, step in enumerate(steps):
    color = ACCENT_BLUE if i < 4 else GREEN_OK
    add_text_box(sl, Inches(0.7), Inches(2.8) + Inches(i * 0.45), Inches(5.5), Inches(0.4),
                 step, font_size=13, color=DARK_TEXT)

# Right: Results
add_rect(sl, Inches(6.8), Inches(2.1), Inches(6.0), Inches(4.5), WHITE)
add_rect(sl, Inches(6.8), Inches(2.1), Inches(6.0), Pt(5), GREEN_OK)
add_text_box(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(0.3),
             "VERIFICATION RESULTS", font_size=16, bold=True, color=GREEN_OK)

results = [
    ("20/20", "Outbox entries created"),
    ("20/20", "Entries replayed successfully"),
    ("0", "Entries stuck in Pending state"),
    ("~2s", "Startup BulkSync completion"),
    ("LWW", "Conflict resolution (last-write-wins)"),
    ("100%", "Data integrity maintained"),
]
for i, (num, desc) in enumerate(results):
    y = Inches(2.8) + Inches(i * 0.6)
    add_rect(sl, Inches(7.0), y, Inches(0.8), Inches(0.35), GREEN_OK)
    add_text_box(sl, Inches(7.0), y, Inches(0.8), Inches(0.35),
                 num, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(7.9), y, Inches(4.5), Inches(0.35),
                 desc, font_size=14, color=DARK_TEXT)

# Bottom banner
add_rect(sl, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.35), DARK_NAVY)
add_text_box(sl, Inches(0.7), Inches(6.8), Inches(11.9), Inches(0.35),
             " The transactional outbox pattern provides reliable, at-least-once delivery semantics with zero data loss.",
             font_size=12, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — CONTRIBUTIONS
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_header(sl, 7, "Contributions & Lessons Learned",
                   "Five technical contributions and key insights from the project")
add_footer(sl)

# Contributions
add_text_box(sl, Inches(0.6), Inches(1.9), Inches(6), Inches(0.35),
             "TECHNICAL CONTRIBUTIONS", font_size=16, bold=True, color=ACCENT_BLUE)

contribs = [
    ("1", "Outbox-Pattern Offline Sync Engine",
     "Adapted the distributed-systems transactional outbox pattern to a single-user SQLite-backed desktop context. Lightweight, zero external dependencies, <1ms overhead per write."),
    ("2", "Dual-Database Architecture",
     "SQLite (primary) + MongoDB (relay) with ISyncableEntity interface exposing ToMongoDocument(). Keeping transformation logic on entities simplified adding new syncable types."),
    ("3", "SSE-Streamed Multi-Modal AI Chatbot",
     "Integrated streaming HTTP responses inside a Tauri-hosted application without WebView buffering issues. SSE request issued directly from React renderer to local ASP.NET Core server."),
    ("4", "Self-Contained Desktop Packaging",
     "Bundled ASP.NET Core server, React frontend, and Tauri shell into a single ~35 MB installer. Stdout-signal protocol decouples frontend load from backend start time."),
    ("5", "Tauri Sidecar + Splashscreen Lifecycle",
     "Two-window lifecycle with Rust backend management. 9 IPC commands bridge shell and web renderer. Neither frontend nor backend required modification for desktop shell integration."),
]

for i, (num, title, desc) in enumerate(contribs):
    y = Inches(2.4) + Inches(i * 0.95)
    add_rect(sl, Inches(0.6), y, Inches(0.4), Inches(0.4), ACCENT_BLUE)
    add_text_box(sl, Inches(0.6), y, Inches(0.4), Inches(0.4),
                 num, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(1.15), y, Inches(5.5), Inches(0.3),
                 title, font_size=14, bold=True, color=DARK_NAVY)
    add_text_box(sl, Inches(1.15), y + Inches(0.3), Inches(5.5), Inches(0.55),
                 desc, font_size=10, color=DARK_TEXT, line_spacing=1.2)

# Lessons Learned
add_text_box(sl, Inches(7.2), Inches(1.9), Inches(5.5), Inches(0.35),
             "LESSONS LEARNED", font_size=16, bold=True, color=TEAL)
lessons = [
    "Offline-first requires deliberate UI design — ConnectivityBar with pending count, colour-coded states, animated progress bar",
    "Choosing conflict resolution early (last-write-wins) shaped entity interfaces and prevented late-stage rework",
    "SignalR group-per-user (Groups.AddToGroupAsync) is a reliable notification pattern that handles reconnection transparently",
    "Clean architectural boundaries enable shell-layer substitution — Rust Tauri shell entirely decoupled from React frontend and .NET backend",
]
for i, lesson in enumerate(lessons):
    y = Inches(2.4) + Inches(i * 0.75)
    add_rect(sl, Inches(7.2), y + Pt(4), Pt(6), Pt(6), TEAL)
    add_text_box(sl, Inches(7.5), y, Inches(5.2), Inches(0.65),
                 lesson, font_size=12, color=DARK_TEXT, line_spacing=1.3)

# Limitations
add_text_box(sl, Inches(7.2), Inches(5.5), Inches(5.5), Inches(0.35),
             "ACKNOWLEDGED LIMITATIONS", font_size=16, bold=True, color=AMBER)
limitations = [
    "• Last-write-wins: no merge dialogue for concurrent edits",
    "• Sub-task backend not implemented (UI scaffold only)",
    "• Mistral API key hardcoded (not via environment variable)",
    "• No automated test suite (manual testing only)",
    "• macOS not notarised (Gatekeeper override required)",
]
for i, lim in enumerate(limitations):
    add_text_box(sl, Inches(7.2), Inches(5.95) + Inches(i * 0.30), Inches(5.5), Inches(0.3),
                 lim, font_size=11, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_header(sl, 8, "Future Work",
                   "Short-term improvements and long-term vision")
add_footer(sl)

# Short-term
add_text_box(sl, Inches(0.6), Inches(1.9), Inches(5.5), Inches(0.4),
             "SHORT-TERM (IMMEDIATE NEXT STEPS)", font_size=16, bold=True, color=ACCENT_BLUE)
add_rect(sl, Inches(0.6), Inches(2.35), Inches(5.5), Pt(2), ACCENT_BLUE)

short_items = [
    ("Sub-task Backend", "Complete the backend support for hierarchical task breakdown"),
    ("API Key Security", "Move Mistral API key to environment variable / secure config"),
    ("Automated Test Suite", "Add unit and integration tests (xUnit, Moq, Testcontainers)"),
    ("Per-Field Conflict Resolution", "Replace last-write-wins with field-level merge"),
    ("Cross-Platform Builds", "Extend installer pipeline to macOS (.dmg) and Linux (.AppImage)"),
]

for i, (title, desc) in enumerate(short_items):
    y = Inches(2.55) + Inches(i * 0.75)
    add_rect(sl, Inches(0.6), y + Pt(3), Inches(0.35), Inches(0.35), ACCENT_BLUE)
    add_text_box(sl, Inches(0.6), y + Pt(3), Inches(0.35), Inches(0.35),
                 str(i+1), font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(1.1), y, Inches(5.0), Inches(0.25),
                 title, font_size=14, bold=True, color=DARK_NAVY)
    add_text_box(sl, Inches(1.1), y + Inches(0.3), Inches(5.0), Inches(0.3),
                 desc, font_size=11, color=DARK_TEXT)

# Long-term
add_text_box(sl, Inches(7.0), Inches(1.9), Inches(5.8), Inches(0.4),
             "LONG-TERM VISION", font_size=16, bold=True, color=TEAL)
add_rect(sl, Inches(7.0), Inches(2.35), Inches(5.8), Pt(2), TEAL)

long_items = [
    ("Mobile Companion App", "Native mobile clients with offline-first sync to the same backend"),
    ("AI Natural-Language Task Creation", "Create tasks by simply describing them in natural language"),
    ("CRDT-Based Sync", "Replace LWW with Conflict-Free Replicated Data Types for true real-time collaboration"),
    ("Analytics Dashboard", "Advanced visualisations: burndown charts, velocity metrics, cycle time analysis"),
    ("Third-Party Plugin System", "Extend via plugins for Git integration, Slack, Teams, Jira import, etc."),
    ("Cloud-Hosted Multi-Tenant SaaS", "Offer TaskFlow as a hosted service with team management and billing"),
]

for i, (title, desc) in enumerate(long_items):
    y = Inches(2.55) + Inches(i * 0.7)
    add_rect(sl, Inches(7.0), y + Pt(3), Inches(0.35), Inches(0.35), TEAL)
    add_text_box(sl, Inches(7.0), y + Pt(3), Inches(0.35), Inches(0.35),
                 str(i+1), font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(7.5), y, Inches(5.3), Inches(0.25),
                 title, font_size=14, bold=True, color=DARK_NAVY)
    add_text_box(sl, Inches(7.5), y + Inches(0.3), Inches(5.3), Inches(0.3),
                 desc, font_size=11, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — DEMO SLIDE
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_rect(sl, 0, 0, SLD_W, SLD_H, DARK_NAVY)
add_rect(sl, 0, 0, Inches(0.15), SLD_H, ACCENT_BLUE)

add_text_box(sl, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.6),
             "Live Demonstration", font_size=40, bold=True, color=WHITE)
add_rect(sl, Inches(0.8), Inches(1.75), Inches(3), Pt(4), ACCENT_BLUE)

add_text_box(sl, Inches(0.8), Inches(2.3), Inches(11.5), Inches(0.4),
             "Three demo scenarios:",
             font_size=20, bold=True, color=ACCENT_BLUE)

demos = [
    ("Scenario 1 — Offline-First Sync",
     "1. Create tasks while online → Observe SQLite + MongoDB sync\n"
     "2. Toggle Manual Offline mode → ConnectivityBar turns amber\n"
     "3. Create more tasks → They persist locally in SQLite\n"
     "4. Reconnect → Watch Outbox replay in real-time → All tasks appear in MongoDB"),
    ("Scenario 2 — AI Assistant",
     "1. Open AI Chatbot → Ask a general question (General mode)\n"
     "2. Switch to Coder mode → Ask to generate a sorting algorithm\n"
     "3. Observe SSE streaming: tokens appear one by one\n"
     "4. Upload a document/ image → Switch to OCR mode → Extract text"),
    ("Scenario 3 — Task Visualisations",
     "1. Open Default view → See temporally grouped tasks\n"
     "2. Switch to Kanban → Drag a task between status columns\n"
     "3. Switch to Gantt → See timeline view\n"
     "4. Open Calendar → See tasks on the weekly grid"),
]

for i, (title, steps) in enumerate(demos):
    y = Inches(2.9) + Inches(i * 1.4)
    add_text_box(sl, Inches(0.8), y, Inches(11.5), Inches(0.3),
                 title, font_size=16, bold=True, color=WHITE)
    add_text_box(sl, Inches(0.8), y + Inches(0.35), Inches(11.5), Inches(1.0),
                 steps, font_size=12, color=LIGHT_GRAY, line_spacing=1.4)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — THANK YOU / Q&A
# ══════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_rect(sl, 0, 0, SLD_W, SLD_H, DARK_NAVY)
add_rect(sl, 0, 0, Inches(0.15), SLD_H, ACCENT_BLUE)

add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.0),
             "Thank You", font_size=52, bold=True, color=WHITE)
add_rect(sl, Inches(0.8), Inches(2.7), Inches(4), Pt(4), ACCENT_BLUE)

add_text_box(sl, Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.5),
             "Questions & Discussion", font_size=28, bold=True, color=ACCENT_BLUE)

add_text_box(sl, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.0),
             "Ahmed Hassan Abdelbarr\n\n"
             "Supervised by: Assoc. Prof. Wael Zakaria Abdallah\n"
             "German University in Cairo\n"
             "Faculty of Engineering and Materials Science\n"
             "Department of Computer Science and Engineering\n\n"
             "2025",
             font_size=18, color=LIGHT_GRAY, line_spacing=1.5)

add_image_safe(sl, "GUC_logo.png", Inches(9.5), Inches(4.5), height=Inches(1.5))


# ══════════════════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"✓ Presentation saved to {OUTPUT}")
print(f"  Slides: {len(prs.slides)}")
print(f"  Dimensions: 16:9 widescreen (13.333 × 7.5 inches)")
